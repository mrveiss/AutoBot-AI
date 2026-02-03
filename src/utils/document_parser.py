# AutoBot - AI-Powered Automation Platform
# Copyright (c) 2025 mrveiss
# Author: mrveiss
"""
Universal Document Parser for Knowledge Base
Supports: PDF, DOC, DOCX, XLSX, PPT, PPTX, ODT, ODS, ODP, ODG
"""

import asyncio
import logging
from pathlib import Path
from typing import Dict, Optional, Tuple

logger = logging.getLogger(__name__)

# Performance optimization: O(1) lookup for file extension checks (Issue #326)
DOCX_EXTENSIONS = {".docx"}
POWERPOINT_EXTENSIONS = {".ppt", ".pptx"}
OPENDOCUMENT_GRAPHICS_EXTENSIONS = {".odp", ".odg"}


class DocumentParser:
    """Universal document parser for extracting text from various formats"""

    def __init__(self):
        """Initialize document parser"""
        self.supported_formats = {
            # PDF
            ".pdf": "application/pdf",
            # Microsoft Office
            ".doc": "application/msword",
            ".docx": (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            ".xlsx": (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ),
            ".ppt": "application/vnd.ms-powerpoint",
            ".pptx": (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            # OpenDocument
            ".odt": "application/vnd.oasis.opendocument.text",
            ".ods": "application/vnd.oasis.opendocument.spreadsheet",
            ".odp": "application/vnd.oasis.opendocument.presentation",
            ".odg": "application/vnd.oasis.opendocument.graphics",
        }

    async def extract_text(
        self, file_path: Path, mime_type: Optional[str] = None
    ) -> Tuple[str, Dict[str, any]]:
        """
        Extract text from document

        Args:
            file_path: Path to document file
            mime_type: Optional MIME type (will be inferred from extension if not provided)

        Returns:
            Tuple of (extracted_text, metadata_dict)
        """
        # Issue #358 - avoid blocking
        exists = await asyncio.to_thread(file_path.exists)
        if not exists:
            raise FileNotFoundError(f"Document not found: {file_path}")

        extension = file_path.suffix.lower()

        if extension not in self.supported_formats:
            raise ValueError(
                f"Unsupported document format: {extension}. "
                f"Supported: {', '.join(self.supported_formats.keys())}"
            )

        # Run extraction in thread pool to avoid blocking
        loop = asyncio.get_event_loop()
        text, metadata = await loop.run_in_executor(
            None, self._extract_text_sync, file_path, extension
        )

        return text, metadata

    def _get_parser_for_extension(self, extension: str):
        """Get parser function for extension (Issue #315 - dispatch table)."""
        # Direct extension mapping for single-extension formats
        single_ext_parsers = {
            ".pdf": self._parse_pdf,
            ".xlsx": self._parse_xlsx,
            ".odt": self._parse_odt,
            ".ods": self._parse_ods,
            ".doc": self._parse_doc_fallback,
        }

        if extension in single_ext_parsers:
            return single_ext_parsers[extension]

        # Check extension sets for multi-extension formats (O(1) lookup)
        if extension in DOCX_EXTENSIONS:
            return self._parse_docx
        if extension in POWERPOINT_EXTENSIONS:
            return self._parse_pptx
        if extension in OPENDOCUMENT_GRAPHICS_EXTENSIONS:
            return self._parse_odp

        return None

    def _extract_text_sync(
        self, file_path: Path, extension: str
    ) -> Tuple[str, Dict[str, any]]:
        """Synchronous text extraction (Issue #315 - refactored to dispatch table)."""
        metadata = {
            "file_name": file_path.name,
            "file_size": file_path.stat().st_size,
            "format": extension,
        }

        try:
            parser = self._get_parser_for_extension(extension)
            if parser is None:
                raise ValueError(f"No parser implemented for {extension}")

            text = parser(file_path, metadata)
            metadata["extraction_success"] = True
            metadata["text_length"] = len(text)

            return text, metadata

        except Exception as e:
            logger.error("Failed to parse %s: %s", file_path, e, exc_info=True)
            metadata["extraction_success"] = False
            metadata["extraction_error"] = str(e)
            return "", metadata

    def _parse_pdf(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from PDF using pypdf"""
        from pypdf import PdfReader

        reader = PdfReader(str(file_path))
        metadata["page_count"] = len(reader.pages)

        text_parts = []
        for page_num, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{text}")
            except Exception as e:
                logger.warning("Failed to extract page %s: %s", page_num, e)

        return "\n\n".join(text_parts)

    def _parse_docx(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from DOCX using python-docx"""
        from docx import Document

        doc = Document(str(file_path))

        # Extract paragraphs
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]

        # Extract text from tables
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)

        metadata["paragraph_count"] = len(paragraphs)
        metadata["table_count"] = len(doc.tables)

        all_text = paragraphs + (["--- Tables ---"] if table_text else []) + table_text
        return "\n".join(all_text)

    def _parse_xlsx(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from Excel using openpyxl"""
        from openpyxl import load_workbook

        wb = load_workbook(str(file_path), data_only=True)
        metadata["sheet_count"] = len(wb.sheetnames)

        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")

            for row in sheet.iter_rows(values_only=True):
                # Filter out empty rows
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(val.strip() for val in row_values):
                    text_parts.append(" | ".join(row_values))

        return "\n".join(text_parts)

    def _parse_pptx(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from PowerPoint using python-pptx"""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        metadata["slide_count"] = len(prs.slides)

        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {slide_num} ---")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)

            # Extract speaker notes if present
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    text_parts.append(f"[Speaker Notes: {notes_text}]")

        return "\n".join(text_parts)

    def _parse_odt(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from ODT using odfpy"""
        from odf import text as odf_text
        from odf.opendocument import load

        doc = load(str(file_path))

        # Extract all text elements
        text_parts = []
        for element in doc.getElementsByType(odf_text.P):
            paragraph_text = self._extract_odf_text(element)
            if paragraph_text.strip():
                text_parts.append(paragraph_text)

        metadata["paragraph_count"] = len(text_parts)
        return "\n".join(text_parts)

    def _parse_ods(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from ODS using odfpy"""
        from odf.opendocument import load
        from odf.table import Table, TableCell, TableRow

        doc = load(str(file_path))

        text_parts = []
        tables = doc.getElementsByType(Table)
        metadata["table_count"] = len(tables)

        for table in tables:
            table_name = table.getAttribute("name")
            text_parts.append(f"--- Table: {table_name} ---")

            for row in table.getElementsByType(TableRow):
                cells = row.getElementsByType(TableCell)
                row_values = [self._extract_odf_text(cell) for cell in cells]
                if any(val.strip() for val in row_values):
                    text_parts.append(" | ".join(row_values))

        return "\n".join(text_parts)

    def _parse_odp(self, file_path: Path, metadata: Dict) -> str:
        """Extract text from ODP/ODG using odfpy"""
        from odf import text as odf_text
        from odf.draw import Page
        from odf.opendocument import load

        doc = load(str(file_path))

        text_parts = []
        pages = doc.getElementsByType(Page)
        metadata["page_count"] = len(pages)

        for page_num, page in enumerate(pages, 1):
            text_parts.append(f"--- Page {page_num} ---")

            # Extract all text from page
            for element in page.getElementsByType(odf_text.P):
                text = self._extract_odf_text(element)
                if text.strip():
                    text_parts.append(text)

        return "\n".join(text_parts)

    def _extract_odf_text(self, element) -> str:
        """Extract text from ODF element recursively"""
        text_parts = []

        # Get text content
        if hasattr(element, "data"):
            text_parts.append(str(element.data))

        # Recursively extract from child elements
        if hasattr(element, "childNodes"):
            for child in element.childNodes:
                child_text = self._extract_odf_text(child)
                if child_text:
                    text_parts.append(child_text)

        return "".join(text_parts)

    def _parse_doc_fallback(self, file_path: Path, metadata: Dict) -> str:
        """
        Fallback parser for old .doc format

        Note: Requires textract or antiword to be installed.
        If not available, returns empty string with warning.
        """
        try:
            import textract

            text = textract.process(str(file_path)).decode("utf-8")
            return text
        except ImportError:
            logger.warning(
                "textract not installed - cannot parse .doc files. "
                "Install with: pip install textract"
            )
            metadata["warning"] = "textract required for .doc files"
            return ""
        except Exception as e:
            logger.error("Failed to parse .doc file: %s", e)
            return ""


# Singleton instance
document_parser = DocumentParser()
