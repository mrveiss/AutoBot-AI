# AutoBot - AI-Powered Automation Platform
# Copyright (c) 2026 mrveiss
# Author: mrveiss
"""
OneDrive / SharePoint Knowledge Connector (Issue #9004)

Indexes documents from OneDrive personal/business and SharePoint document
libraries into the AutoBot knowledge base via Microsoft Graph API.

Supported file types:
- Word documents (.docx) - full text extraction
- Excel spreadsheets (.xlsx) - sheet content as text
- PDF documents (.pdf) - text extraction
- PowerPoint presentations (.pptx) - slide text
- Markdown files (.md) - raw content

Config keys (under ``ConnectorConfig.config``):
    token (str): Microsoft Graph API access token (OAuth2). Required.
    source_type (str): "onedrive" or "sharepoint". Default "onedrive".
    drive_id (str): Specific drive ID to sync. Optional for OneDrive (uses default),
                    required for SharePoint document libraries.
    site_id (str): SharePoint site ID. Required when source_type is "sharepoint".
    folder_path (str): Specific folder path to sync. Default "/" (root).
    sync_subfolders (bool): Recursively sync subfolders. Default True.
    max_file_size_mb (int): Skip files larger than this. Default 100MB.
    supported_extensions (list[str]): File extensions to index.
                                      Default [".docx", ".xlsx", ".pdf", ".pptx", ".md"].
"""

import hashlib
import io
from datetime import datetime
from typing import Any, Dict, List, Optional
from urllib.parse import quote

import aiohttp

from autobot_shared.auth import BearerAuth
from autobot_shared.logging_manager import get_logger
from autobot_shared.time_utils import now_utc, parse_utc_iso
from knowledge.connectors.base import AbstractConnector, RetryableError
from knowledge.connectors.models import (
    ChangeInfo,
    ConnectorConfig,
    ContentResult,
    SourceInfo,
)
from knowledge.connectors.registry import ConnectorRegistry

logger = get_logger(__name__)

_GRAPH_API_BASE = "https://graph.microsoft.com/v1.0"
_REDIS_TS_PREFIX = "connector:onedrive:ts:"
_REDIS_TS_TTL = 86400 * 30  # 30 days

# Default supported extensions
_DEFAULT_EXTENSIONS = [".docx", ".xlsx", ".pdf", ".pptx", ".md", ".txt"]

# Max file size in bytes (default 100MB)
_DEFAULT_MAX_FILE_SIZE = 100 * 1024 * 1024


def _content_hash(text: str) -> str:
    """Generate SHA-256 hash of content for change detection."""
    return hashlib.sha256(text.encode("utf-8")).hexdigest()


def _extract_text_from_docx(content_bytes: bytes) -> str:
    """Extract text from Word .docx file."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content_bytes))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception as exc:
        logger.warning("Failed to extract text from DOCX: %s", exc)
        return ""


def _extract_text_from_xlsx(content_bytes: bytes) -> str:
    """Extract text from Excel .xlsx file (all sheets as text)."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content_bytes), read_only=True, data_only=True)
        sheets_text = []
        for sheet in wb:
            sheet_lines = [f"## Sheet: {sheet.title}"]
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    sheet_lines.append(row_text)
            sheets_text.append("\n".join(sheet_lines))
        return "\n\n".join(sheets_text)
    except Exception as exc:
        logger.warning("Failed to extract text from XLSX: %s", exc)
        return ""


def _extract_text_from_pdf(content_bytes: bytes) -> str:
    """Extract text from PDF file."""
    try:
        from PyPDF2 import PdfReader

        pdf = PdfReader(io.BytesIO(content_bytes))
        pages_text = []
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text()
            if text.strip():
                pages_text.append(f"## Page {page_num}\n{text}")
        return "\n\n".join(pages_text)
    except Exception as exc:
        logger.warning("Failed to extract text from PDF: %s", exc)
        return ""


def _extract_text_from_pptx(content_bytes: bytes) -> str:
    """Extract text from PowerPoint .pptx file."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content_bytes))
        slides_text = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines = [f"## Slide {slide_num}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text)
            if len(slide_lines) > 1:  # Only add if there's text
                slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)
    except Exception as exc:
        logger.warning("Failed to extract text from PPTX: %s", exc)
        return ""


async def _load_ts(connector_id: str, source_id: str) -> Optional[str]:
    """Load last-modified timestamp for a source from Redis."""
    try:
        from autobot_shared.redis_client import get_redis_client

        redis = get_redis_client(database="knowledge")
        key = f"{_REDIS_TS_PREFIX}{connector_id}:{source_id}"
        value = redis.get(key)
        if hasattr(value, "__await__"):
            value = await value
        if isinstance(value, bytes):
            return value.decode("utf-8")
        return value
    except Exception as exc:
        logger.warning("Redis load_ts failed for %s: %s", source_id, exc)
        return None


async def _store_ts(connector_id: str, source_id: str, ts: str) -> None:
    """Store last-modified timestamp for a source in Redis."""
    try:
        from autobot_shared.redis_client import get_redis_client

        redis = get_redis_client(database="knowledge")
        key = f"{_REDIS_TS_PREFIX}{connector_id}:{source_id}"
        result = redis.set(key, ts, ex=_REDIS_TS_TTL)
        if hasattr(result, "__await__"):
            await result
    except Exception as exc:
        logger.warning("Redis store_ts failed for %s: %s", source_id, exc)


@ConnectorRegistry.register("onedrive")
class OneDriveConnector(AbstractConnector):
    """Knowledge connector for OneDrive and SharePoint document libraries.

    Indexes documents from OneDrive personal/business or SharePoint libraries
    into the AutoBot knowledge base. Supports incremental sync via last-modified
    timestamp comparison.

    Each file becomes one KB fact keyed by ``onedrive:{connector_id}:file:{file_id}``.
    """

    connector_type = "onedrive"
    # Issue #4421: needs OAuth2 access token (tier 2 = credentials/OAuth)
    tier = 2

    @classmethod
    def auth_schema(cls) -> type:
        """OneDrive requires a bearer token (OAuth2 access token) — Issue #8145."""
        return BearerAuth

    @classmethod
    def output_schema(cls) -> Dict[str, Any]:
        """Return JSONSchema for ContentResult.metadata (Issue #8147)."""
        return {
            "type": "object",
            "required": ["file_id", "file_name", "file_extension", "last_modified"],
            "properties": {
                "file_id": {"type": "string", "description": "OneDrive file ID"},
                "file_name": {"type": "string", "description": "File name with extension"},
                "file_extension": {"type": "string", "description": "File extension (e.g., .docx)"},
                "file_path": {"type": "string", "description": "Full path in drive"},
                "file_size": {"type": "integer", "description": "File size in bytes"},
                "last_modified": {"type": "string", "description": "ISO-8601 last modified datetime"},
                "web_url": {"type": "string", "description": "Web URL to view file"},
                "drive_id": {"type": "string", "description": "Drive ID"},
                "site_id": {"type": "string", "description": "SharePoint site ID (if applicable)"},
            },
        }

    @property
    def max_concurrency(self) -> int:
        """Max parallel file fetches (Issue #8148). Default 4 for OneDrive."""
        cfg_val = self.config.max_concurrency
        return cfg_val if cfg_val is not None else 4

    def __init__(self, config: ConnectorConfig) -> None:
        super().__init__(config)
        cfg = config.config
        self._token: str = cfg.get("token", "")
        self._source_type: str = cfg.get("source_type", "onedrive")  # "onedrive" or "sharepoint"
        self._drive_id: Optional[str] = cfg.get("drive_id")
        self._site_id: Optional[str] = cfg.get("site_id")
        self._folder_path: str = cfg.get("folder_path", "/")
        self._sync_subfolders: bool = cfg.get("sync_subfolders", True)
        self._max_file_size: int = cfg.get("max_file_size_mb", 100) * 1024 * 1024
        self._supported_extensions: List[str] = cfg.get("supported_extensions", _DEFAULT_EXTENSIONS)
        self._graph_url: str = cfg.get("graph_api_base", _GRAPH_API_BASE).rstrip("/")

    # ------------------------------------------------------------------
    # AbstractConnector interface
    # ------------------------------------------------------------------

    async def test_connection(self) -> bool:
        """Verify OneDrive/SharePoint credentials by fetching drive metadata."""
        try:
            drive_url = await self._get_drive_url()
            result = await self._graph_request("GET", drive_url)
            healthy = result.get("status_code") == 200
            if not healthy:
                self.logger.warning(
                    "OneDrive test_connection failed: HTTP %s",
                    result.get("status_code"),
                )
            return healthy
        except Exception as exc:
            self.logger.error("OneDrive test_connection raised: %s", exc)
            return False

    async def discover_sources(self) -> List[SourceInfo]:
        """Return a SourceInfo for every supported file in the configured drive/folder."""
        files = await self._list_all_files()
        sources: List[SourceInfo] = []
        for file_item in files:
            source_info = self._file_to_source_info(file_item)
            if source_info:
                sources.append(source_info)
        return sources

    async def fetch_content(self, source_id: str) -> ContentResult | None:
        """Fetch and extract text content for a OneDrive file by ID."""
        # source_id format: "onedrive:{connector_id}:file:{file_id}"
        parts = source_id.split(":")
        if len(parts) != 4 or parts[0] != "onedrive" or parts[2] != "file":
            self.logger.error("Invalid source_id format: %s", source_id)
            return None

        file_id = parts[3]

        # Get file metadata first
        drive_url = await self._get_drive_url()
        metadata_url = f"{drive_url}/items/{file_id}"
        metadata_result = await self._graph_request("GET", metadata_url)

        if metadata_result.get("status_code") != 200:
            self.logger.error(
                "Failed to fetch OneDrive file metadata %s: HTTP %s",
                file_id,
                metadata_result.get("status_code"),
            )
            return None

        file_meta = metadata_result.get("body", {})
        file_name = file_meta.get("name", "")
        file_ext = self._get_file_extension(file_name).lower()

        # Download file content
        download_url = f"{metadata_url}/content"
        download_result = await self._graph_request("GET", download_url, raw_content=True)

        if download_result.get("status_code") != 200:
            self.logger.error(
                "Failed to download OneDrive file %s: HTTP %s",
                file_id,
                download_result.get("status_code"),
            )
            return None

        content_bytes = download_result.get("content", b"")

        # Extract text based on file type
        text = ""
        if file_ext == ".docx":
            text = _extract_text_from_docx(content_bytes)
        elif file_ext == ".xlsx":
            text = _extract_text_from_xlsx(content_bytes)
        elif file_ext == ".pdf":
            text = _extract_text_from_pdf(content_bytes)
        elif file_ext == ".pptx":
            text = _extract_text_from_pptx(content_bytes)
        elif file_ext in [".md", ".txt"]:
            try:
                text = content_bytes.decode("utf-8")
            except UnicodeDecodeError:
                self.logger.warning("Failed to decode text file %s as UTF-8", file_id)
                text = content_bytes.decode("utf-8", errors="replace")
        else:
            self.logger.warning("Unsupported file extension for %s: %s", file_id, file_ext)
            return None

        if not text.strip():
            self.logger.debug("Extracted empty text from file %s", file_id)
            return None

        # Add file header
        header = f"# {file_name}\n\nFile: {file_meta.get('webUrl', '')}\n\n"
        text = header + text

        # Update stored timestamp
        last_modified = file_meta.get("lastModifiedDateTime", "")
        if last_modified:
            await _store_ts(self.config.connector_id, source_id, last_modified)

        return ContentResult(
            source_id=source_id,
            content=text,
            content_type="text/plain",
            metadata={
                "file_id": file_id,
                "file_name": file_name,
                "file_extension": file_ext,
                "file_path": file_meta.get("parentReference", {}).get("path", ""),
                "file_size": file_meta.get("size", 0),
                "last_modified": last_modified,
                "web_url": file_meta.get("webUrl", ""),
                "drive_id": file_meta.get("parentReference", {}).get("driveId", ""),
                "site_id": self._site_id or "",
                "connector_id": self.config.connector_id,
            },
        )

    async def detect_changes(self, since: datetime | None = None) -> List[ChangeInfo]:
        """Return ChangeInfo for files added or modified since *since*.

        When *since* is None all files are reported as 'added'. Otherwise
        compares lastModifiedDateTime against stored Redis timestamp.
        """
        files = await self._list_all_files()
        changes: List[ChangeInfo] = []

        for file_item in files:
            file_id = file_item.get("id", "")
            source_id = f"onedrive:{self.config.connector_id}:file:{file_id}"
            last_modified = file_item.get("lastModifiedDateTime", "")

            change = await self._classify_change(source_id, last_modified, since)
            if change:
                changes.append(change)

        return changes

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    async def _get_drive_url(self) -> str:
        """Build the Graph API drive endpoint URL based on configuration."""
        if self._source_type == "sharepoint":
            if not self._site_id:
                raise ValueError("site_id required when source_type is 'sharepoint'")
            if self._drive_id:
                return f"{self._graph_url}/sites/{self._site_id}/drives/{self._drive_id}"
            else:
                # Use default document library
                return f"{self._graph_url}/sites/{self._site_id}/drive"
        else:  # onedrive
            if self._drive_id:
                return f"{self._graph_url}/drives/{self._drive_id}"
            else:
                # Use current user's default drive
                return f"{self._graph_url}/me/drive"

    async def _list_all_files(self) -> List[Dict[str, Any]]:
        """List all files in the configured drive/folder (with pagination)."""
        drive_url = await self._get_drive_url()

        # Build folder item URL
        if self._folder_path == "/":
            folder_url = f"{drive_url}/root/children"
        else:
            # Encode path for URL
            encoded_path = quote(self._folder_path.strip("/"), safe="")
            folder_url = f"{drive_url}/root:/{encoded_path}:/children"

        files: List[Dict[str, Any]] = []
        next_url = folder_url

        while next_url:
            result = await self._graph_request("GET", next_url)
            if result.get("status_code") != 200:
                self.logger.error(
                    "Failed to list OneDrive files: HTTP %s",
                    result.get("status_code"),
                )
                break

            body = result.get("body", {})
            items = body.get("value", [])

            for item in items:
                if "file" in item:  # It's a file, not a folder
                    # Check file size
                    size = item.get("size", 0)
                    if size > self._max_file_size:
                        self.logger.debug(
                            "Skipping file %s (size %d exceeds max %d)",
                            item.get("name"),
                            size,
                            self._max_file_size,
                        )
                        continue

                    # Check file extension
                    file_name = item.get("name", "")
                    file_ext = self._get_file_extension(file_name).lower()
                    if file_ext in self._supported_extensions:
                        files.append(item)
                    else:
                        self.logger.debug(
                            "Skipping file %s (extension %s not supported)",
                            file_name,
                            file_ext,
                        )
                elif "folder" in item and self._sync_subfolders:
                    # Recursively list subfolder
                    subfolder_files = await self._list_folder_recursive(item.get("id", ""))
                    files.extend(subfolder_files)

            # Handle pagination
            next_url = body.get("@odata.nextLink")

        return files

    async def _list_folder_recursive(self, folder_id: str) -> List[Dict[str, Any]]:
        """Recursively list all files in a folder by ID."""
        drive_url = await self._get_drive_url()
        folder_url = f"{drive_url}/items/{folder_id}/children"
        files: List[Dict[str, Any]] = []
        next_url = folder_url

        while next_url:
            result = await self._graph_request("GET", next_url)
            if result.get("status_code") != 200:
                self.logger.warning(
                    "Failed to list folder %s: HTTP %s",
                    folder_id,
                    result.get("status_code"),
                )
                break

            body = result.get("body", {})
            items = body.get("value", [])

            for item in items:
                if "file" in item:
                    size = item.get("size", 0)
                    if size > self._max_file_size:
                        continue

                    file_name = item.get("name", "")
                    file_ext = self._get_file_extension(file_name).lower()
                    if file_ext in self._supported_extensions:
                        files.append(item)
                elif "folder" in item and self._sync_subfolders:
                    subfolder_files = await self._list_folder_recursive(item.get("id", ""))
                    files.extend(subfolder_files)

            next_url = body.get("@odata.nextLink")

        return files

    async def _classify_change(
        self,
        source_id: str,
        last_modified: str,
        since: datetime | None,
    ) -> ChangeInfo | None:
        """Return ChangeInfo when the file is new or was modified after *since*."""
        if since is None:
            return ChangeInfo(
                source_id=source_id,
                change_type="added",
                timestamp=now_utc(),
                details={"last_modified": last_modified},
            )

        stored_ts = await _load_ts(self.config.connector_id, source_id)
        if stored_ts is None or last_modified > stored_ts:
            change_type = "added" if stored_ts is None else "modified"
            return ChangeInfo(
                source_id=source_id,
                change_type=change_type,
                timestamp=parse_utc_iso(last_modified) if last_modified else now_utc(),
                details={"last_modified": last_modified},
            )

        return None

    def _file_to_source_info(self, file_item: Dict[str, Any]) -> SourceInfo | None:
        """Convert a Graph API file item to SourceInfo."""
        file_id = file_item.get("id", "")
        if not file_id:
            return None

        file_name = file_item.get("name", "")
        file_ext = self._get_file_extension(file_name).lower()

        if file_ext not in self._supported_extensions:
            return None

        source_id = f"onedrive:{self.config.connector_id}:file:{file_id}"
        last_modified_str = file_item.get("lastModifiedDateTime", "")
        last_modified = parse_utc_iso(last_modified_str) if last_modified_str else now_utc()

        # Build path from parent reference
        parent_ref = file_item.get("parentReference", {})
        parent_path = parent_ref.get("path", "")
        # Extract the actual folder path from the full path format
        # e.g., "/drive/root:/Documents" -> "/Documents"
        if "/root:" in parent_path:
            folder_path = parent_path.split("/root:", 1)[1]
        else:
            folder_path = "/"
        full_path = f"{folder_path}/{file_name}".replace("//", "/")

        return SourceInfo(
            source_id=source_id,
            name=file_name,
            path=full_path,
            content_type=self._get_content_type(file_ext),
            size_bytes=file_item.get("size", 0),
            last_modified=last_modified,
            metadata={
                "file_id": file_id,
                "file_name": file_name,
                "file_extension": file_ext,
                "web_url": file_item.get("webUrl", ""),
                "drive_id": parent_ref.get("driveId", ""),
                "site_id": self._site_id or "",
            },
        )

    @staticmethod
    def _get_file_extension(file_name: str) -> str:
        """Extract file extension from file name."""
        if "." in file_name:
            return "." + file_name.rsplit(".", 1)[-1]
        return ""

    @staticmethod
    def _get_content_type(file_ext: str) -> str:
        """Map file extension to MIME type."""
        content_type_map = {
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".pdf": "application/pdf",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".md": "text/markdown",
            ".txt": "text/plain",
        }
        return content_type_map.get(file_ext.lower(), "application/octet-stream")

    async def _graph_request(
        self,
        method: str,
        url: str,
        json_data: Dict[str, Any] | None = None,
        params: Dict[str, Any] | None = None,
        raw_content: bool = False,
    ) -> Dict[str, Any]:
        """Make an HTTP request to Microsoft Graph API with auth and error handling.

        Args:
            method: HTTP method (GET, POST, etc.)
            url: Full URL to request
            json_data: JSON body for POST/PATCH requests
            params: Query parameters
            raw_content: If True, return raw bytes in "content" key instead of parsing JSON

        Returns:
            Dict with keys: status_code, body (or content if raw_content=True), error
        """
        headers = {
            "Authorization": f"Bearer {self._token}",
            "Content-Type": "application/json",
        }

        try:
            timeout = aiohttp.ClientTimeout(total=60.0)  # Longer timeout for file downloads
            async with aiohttp.ClientSession(timeout=timeout) as session:
                async with session.request(
                    method,
                    url,
                    headers=headers,
                    json=json_data,
                    params=params,
                ) as resp:
                    status_code = resp.status

                    # Handle redirects for download URLs
                    if status_code in (301, 302, 303, 307, 308):
                        location = resp.headers.get("Location")
                        if location and raw_content:
                            # Follow redirect for file download
                            async with session.get(location) as redirect_resp:
                                content = await redirect_resp.read()
                                return {
                                    "status_code": redirect_resp.status,
                                    "content": content,
                                }

                    # Handle no-content responses
                    if status_code == 204:
                        return {"status_code": 204, "body": {}}

                    if raw_content:
                        content = await resp.read()
                        return {
                            "status_code": status_code,
                            "content": content,
                        }

                    try:
                        body = await resp.json()
                    except aiohttp.ContentTypeError:
                        body = {}

                    # Log errors
                    if status_code >= 400:
                        error_msg = body.get("error", {}).get("message", "Unknown error")
                        self.logger.warning(
                            "Graph API request to %s failed: HTTP %d - %s",
                            url,
                            status_code,
                            error_msg,
                        )

                    return {
                        "status_code": status_code,
                        "body": body,
                        "error": body.get("error", {}).get("message") if status_code >= 400 else None,
                    }

        except aiohttp.ClientError as exc:
            self.logger.error("Graph API request to %s failed: %s", url, exc)
            return {
                "status_code": 0,
                "body": {},
                "error": str(exc),
            }
        except Exception as exc:
            self.logger.error("Unexpected error in Graph API request to %s: %s", url, exc)
            return {
                "status_code": 0,
                "body": {},
                "error": str(exc),
            }
